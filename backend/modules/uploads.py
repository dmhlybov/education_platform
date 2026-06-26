import xml.etree.ElementTree as ET
from io import BytesIO
from pathlib import Path
from zipfile import BadZipFile, ZipFile

import aiofiles
from fastapi import APIRouter, Depends, File, HTTPException, Query, UploadFile

from backend.config import SCORM_DIR, UPLOADS_DIR
from backend.core.app.deps import current_active_user_ctx, current_superuser_ctx
from backend.schemas.auth import UserContext


def _unique_dir(parent: Path, stem: str) -> Path:
    candidate = parent / stem
    counter = 0
    while candidate.exists():
        counter += 1
        candidate = parent / f"{stem}-{counter}"
    return candidate


router = APIRouter(prefix="/uploads", tags=["Документы"])

CHUNK_SIZE = 1024 * 1024  # 1 MB


@router.post("/document")
async def upload_document(
    file: UploadFile = File(...),
    user: UserContext = Depends(current_superuser_ctx),
):
    UPLOADS_DIR.mkdir(parents=True, exist_ok=True)

    original_filename = Path(file.filename or "document").name

    original_path = UPLOADS_DIR / original_filename
    stem = original_path.stem
    suffix = original_path.suffix

    counter = 0

    try:
        while True:
            if counter == 0:
                target = UPLOADS_DIR / f"{stem}{suffix}"
            else:
                target = UPLOADS_DIR / f"{stem}-{counter}{suffix}"

            try:
                async with aiofiles.open(target, "xb") as out_file:
                    while chunk := await file.read(CHUNK_SIZE):
                        await out_file.write(chunk)

                break

            except FileExistsError:
                counter += 1

    finally:
        await file.close()

    return {
        "url": f"/uploads/{target.name}",
        "filename": target.name,
    }


@router.get("/pptx-preview")
async def pptx_preview(
    file: str = Query(..., description="Path like /uploads/file.pptx"),
    _user: UserContext = Depends(current_active_user_ctx),
):
    try:
        from pptx import Presentation
    except ImportError:
        raise HTTPException(status_code=501, detail="python-pptx not installed")

    filename = Path(file).name
    file_path = UPLOADS_DIR / filename

    if not file_path.exists() or file_path.suffix.lower() != ".pptx":
        raise HTTPException(status_code=404, detail="PPTX file not found")

    try:
        prs = Presentation(str(file_path))
    except Exception as e:
        raise HTTPException(status_code=422, detail=f"Cannot open PPTX: {e}")

    slides = []
    for i, slide in enumerate(prs.slides):
        title_text = None
        body_texts: list[str] = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            ph = getattr(shape, "placeholder_format", None)
            if ph is not None and ph.idx == 0:
                title_text = text
            else:
                body_texts.append(text)

        slides.append(
            {
                "number": i + 1,
                "title": title_text,
                "texts": body_texts,
            }
        )

    return {"slides": slides, "total": len(slides)}


@router.post("/scorm")
async def upload_scorm(
    file: UploadFile = File(...),
    _user: UserContext = Depends(current_superuser_ctx),
):
    SCORM_DIR.mkdir(parents=True, exist_ok=True)

    stem = Path(file.filename or "scorm").stem
    pkg_dir = _unique_dir(SCORM_DIR, stem)
    pkg_dir.mkdir(parents=True, exist_ok=True)

    content = await file.read()
    await file.close()

    try:
        with ZipFile(BytesIO(content)) as zf:
            pkg_root = str(pkg_dir.resolve())
            for member in zf.namelist():
                dest = (pkg_dir / member).resolve()
                if not str(dest).startswith(pkg_root):
                    raise HTTPException(
                        status_code=422, detail="Invalid ZIP: path traversal"
                    )
            zf.extractall(pkg_dir)
    except BadZipFile:
        import shutil

        shutil.rmtree(pkg_dir, ignore_errors=True)
        raise HTTPException(status_code=422, detail="Not a valid ZIP file")

    manifest_path = pkg_dir / "imsmanifest.xml"
    entry_href = "index.html"
    title = stem

    if manifest_path.exists():
        try:
            root = ET.parse(manifest_path).getroot()
            ns = root.tag.split("}")[0].lstrip("{") if "}" in root.tag else ""
            p = f"{{{ns}}}" if ns else ""

            for elem in root.iter(f"{p}title"):
                t = (elem.text or "").strip()
                if t:
                    title = t
                    break

            for resource in root.iter(f"{p}resource"):
                href = resource.get("href", "").strip()
                if href:
                    entry_href = href
                    break
        except ET.ParseError:
            pass

    pkg_name = pkg_dir.name
    return {
        "url": f"/scorm/{pkg_name}/{entry_href}",
        "title": title,
        "package": pkg_name,
    }
